"""
מייצר מצגת PowerPoint: "כמעט אירוע איכות"
הרצה: python generate_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ─── צבעים ───
C_BLUE   = RGBColor(0x1F, 0x38, 0x64)   # כחול כהה – בטיחות
C_GREEN  = RGBColor(0x1E, 0x56, 0x31)   # ירוק כהה – איכות
C_AMBER  = RGBColor(0xBF, 0x6E, 0x0E)   # ענבר – כותרת שקף 2
C_LGRAY  = RGBColor(0xF0, 0xF0, 0xF0)   # אפור בהיר – תיבת הגדרה
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_GOLD   = RGBColor(0xFF, 0xD7, 0x00)
C_LGREEN = RGBColor(0x90, 0xEE, 0x90)

STEP_COLORS = [
    RGBColor(0x2E, 0x75, 0xB6),  # כחול  – קבלה
    RGBColor(0xED, 0x73, 0x00),  # כתום  – חיתוך
    RGBColor(0xC0, 0x00, 0x00),  # אדום  – גילוי
    RGBColor(0x70, 0x00, 0x20),  # בורדו – בלימה
    RGBColor(0x1E, 0x56, 0x31),  # ירוק  – תיקון
]


# ─── עזרים ───

def rtl(p):
    p._p.get_or_add_pPr().set(qn('a:rtl'), '1')


def colored_box(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def transparent_box(slide, l, t, w, h):
    s = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    return s


def para(tf, text, size, bold=False, color=C_BLACK,
         align=PP_ALIGN.RIGHT, first=False, space_before=0, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    rtl(p)
    if space_before:
        p.space_before = Pt(space_before)
    if text:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return p


def setup_tf(shape, margin_lr=0.15, margin_top=0.12):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(margin_lr)
    tf.margin_right = Inches(margin_lr)
    tf.margin_top   = Inches(margin_top)
    return tf


# ─────────────────────────────────────────────────────────
# שקף 1 — הרעיון
# ─────────────────────────────────────────────────────────

def slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # כותרת ראשית – רקע כחול
    hdr = colored_box(slide, 0, 0, 13.333, 1.15, C_BLUE)
    tf = setup_tf(hdr, margin_top=0.18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    rtl(p)
    r = p.add_run()
    r.text = "כמעט אירוע איכות — מודל פרואקטיבי לזיהוי מוקדם"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = C_WHITE

    # עמודה שמאל — בטיחות
    safety = colored_box(slide, 0.3, 1.3, 5.5, 4.05, C_BLUE)
    tf = setup_tf(safety)
    para(tf, "🛡️  עולם הבטיחות", 18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, first=True)
    para(tf, '"כמעט ונפגע"', 16, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER, space_before=6)
    para(tf, "", 4)
    para(tf, "• אירוע שיכל להסתיים בנפגע —", 13, color=C_WHITE)
    para(tf, "   זוהה לפני שקרה הנזק", 13, color=C_WHITE)
    para(tf, "• כלי פרואקטיבי ומניעתי", 13, color=C_WHITE, space_before=5)
    para(tf, "• הוכח: מקטין תאונות, מעלה מודעות", 13, color=C_WHITE, space_before=5)
    para(tf, "", 5)
    para(tf, "📍 ניסיון אישי: פירוטכניקה", 12, bold=True, color=C_GOLD)

    # חץ מרכז
    arr = transparent_box(slide, 5.97, 2.55, 1.4, 1.8)
    tf = arr.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    rtl(p)
    r = p.add_run()
    r.text = "→"
    r.font.size = Pt(38)
    r.font.color.rgb = RGBColor(0x90, 0x90, 0x90)
    para(tf, "אותו מודל", 12, bold=True, color=RGBColor(0x50, 0x50, 0x50), align=PP_ALIGN.CENTER)
    para(tf, "עולם אחר", 12, bold=True, color=RGBColor(0x50, 0x50, 0x50), align=PP_ALIGN.CENTER)

    # עמודה ימין — איכות
    quality = colored_box(slide, 7.53, 1.3, 5.5, 4.05, C_GREEN)
    tf = setup_tf(quality)
    para(tf, "⚙️  עולם האיכות", 18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, first=True)
    para(tf, '"כמעט אירוע איכות"', 16, bold=True, color=C_LGREEN, align=PP_ALIGN.CENTER, space_before=6)
    para(tf, "", 4)
    para(tf, "• אירוע שזוהה לפני שהפך לחריגה,", 13, color=C_WHITE)
    para(tf, "   פסילה או אירוע לקוח", 13, color=C_WHITE)
    para(tf, "• זיהוי מוקדם — לא רק בתחקיר", 13, color=C_WHITE, space_before=5)
    para(tf, "• תרבות איכות פרואקטיבית", 13, color=C_WHITE, space_before=5)
    para(tf, "", 5)
    para(tf, "📍 יעד: כל מחלקות האיכות", 12, bold=True, color=C_LGREEN)

    # תיבת הגדרה תחתונה
    defbox = colored_box(slide, 0.3, 5.52, 12.73, 1.8, C_LGRAY)
    tf = setup_tf(defbox, margin_lr=0.25, margin_top=0.14)
    para(tf, '📋  הגדרה: "כמעט אירוע איכות"', 14, bold=True,
         color=C_BLUE, align=PP_ALIGN.RIGHT, first=True)
    para(tf,
         "אירוע שבו זוהה פוטנציאל לכשל, אי-התאמה או חריגה בתהליך, במוצר, בחומר או בבקרה —",
         12.5, color=C_BLACK, align=PP_ALIGN.RIGHT, space_before=4)
    para(tf,
         "אשר זוהה ונבלם לפני שהתפתח לאירוע איכות בפועל.",
         12.5, bold=True, color=C_GREEN, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────
# שקף 2 — מקרה לדוגמה
# ─────────────────────────────────────────────────────────

def slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # כותרת
    hdr = colored_box(slide, 0, 0, 13.333, 1.15, C_AMBER)
    tf = setup_tf(hdr, margin_top=0.18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    rtl(p)
    r = p.add_run()
    r.text = "מקרה לדוגמה — כמעט אירוע איכות: מנת גרפיט"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = C_WHITE

    # 5 תחנות ציר זמן
    steps = [
        ("📦", "קבלה", ["2 מנות גרפיט שונות", "סימון על הארגז בלבד"]),
        ("✂️",  "חיתוך", ["שתי המנות נחתכו", "במחסן המתכות"]),
        ("❓",  "גילוי", ["בסיור: אין כיתוב", "על גוף הפריט"]),
        ("🛑",  "בלימה", ["זוהה ונבלם לפני", "שימוש בחומר לא ידוע"]),
        ("✅",  "תיקון", ["נוהל חדש: כיתוב", "על גוף הפריט + פדחת"]),
    ]

    box_w   = 2.35
    arrow_w = 0.2
    start_x = 0.37
    box_y   = 1.35
    box_h   = 4.0

    for i, (icon, title, lines) in enumerate(steps):
        x = start_x + i * (box_w + arrow_w)

        box = colored_box(slide, x, box_y, box_w, box_h, STEP_COLORS[i])
        tf = setup_tf(box, margin_lr=0.1, margin_top=0.18)

        # מספר שלב
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        rtl(p)
        r = p.add_run()
        r.text = str(i + 1)
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.color.rgb = C_WHITE

        para(tf, icon, 26, align=PP_ALIGN.CENTER, color=C_WHITE, space_before=2)
        para(tf, title, 15, bold=True, align=PP_ALIGN.CENTER, color=C_WHITE, space_before=8)
        for line in lines:
            para(tf, line, 12, align=PP_ALIGN.CENTER, color=C_WHITE, space_before=4)

        # חץ בין תחנות
        if i < 4:
            arrow_x = x + box_w + 0.01
            tb = transparent_box(slide, arrow_x, box_y + 1.55, arrow_w, 0.6)
            tf_a = tb.text_frame
            p_a = tf_a.paragraphs[0]
            p_a.alignment = PP_ALIGN.CENTER
            rtl(p_a)
            r_a = p_a.add_run()
            r_a.text = "▶"
            r_a.font.size = Pt(16)
            r_a.font.color.rgb = RGBColor(0xA0, 0xA0, 0xA0)

    # תיבת תוצאה תחתונה
    res = colored_box(slide, 0.37, 5.52, 12.59, 1.8, C_LGRAY)
    tf = setup_tf(res, margin_lr=0.25, margin_top=0.14)
    para(tf, "✅  תוצאה ולקחים", 14, bold=True, color=C_GREEN,
         align=PP_ALIGN.RIGHT, first=True)
    para(tf,
         "שיחת חידוד ויישור קו — נקבע נוהל: כיתוב על גוף הפריט ובפדחת לפני חיתוך",
         12.5, color=C_BLACK, align=PP_ALIGN.RIGHT, space_before=4)
    para(tf,
         "מאז: אפס אירועי ערבוב חומרים  |  כמעט אירוע איכות שנמנע ברגע האחרון",
         12.5, bold=True, color=C_BLUE, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────
# הרצה ראשית
# ─────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide1(prs)
    slide2(prs)

    out = "כמעט_אירוע_איכות.pptx"
    prs.save(out)
    print(f"✅ נשמר: {out}")


if __name__ == "__main__":
    main()
